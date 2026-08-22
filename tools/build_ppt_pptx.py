# -*- coding: utf-8 -*-
"""Generate the technical-demo PPT with python-pptx (layout-optimized).

Plain engineering style: white/light-gray background, dark-blue titles,
card-based grouping, shape-drawn architecture/flow diagrams, and
ready-to-use mermaid code blocks on diagram pages.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

OUT = r"提交材料\AI能力测评智能体_技术演示.pptx"

PRIMARY = RGBColor(0x31, 0x57, 0xD5)
PRIMARY_SOFT = RGBColor(0xE9, 0xEE, 0xFB)
DARK = RGBColor(0x22, 0x2A, 0x3A)
MUTED = RGBColor(0x6A, 0x74, 0x86)
CARD_BG = RGBColor(0xF6, 0xF8, 0xFB)
CARD_LINE = RGBColor(0xDF, 0xE4, 0xEE)
CODE_BG = RGBColor(0xF2, 0xF4, 0xF8)
CODE_LINE = RGBColor(0xD0, 0xD6, 0xE2)
CODE_TEXT = RGBColor(0x3A, 0x4A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
EA_FONT = "微软雅黑"

EMU_W = 12192000
EMU_H = 6858000


def set_ea(run, name=EA_FONT):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = rPr.makeelement(qn("a:ea"), {})
        rPr.append(ea)
    ea.set("typeface", name)


def add_text(slide, x, y, w, h, lines, size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.15, space_after=6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    box.fill.background()
    box.line.fill.background()
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        set_ea(run)
    return box


def add_shape(slide, shape_type, x, y, w, h, text="", fill=CARD_BG, line=CARD_LINE, size=12,
              color=DARK, bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    sp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    sp.shadow.inherit = False
    if text:
        tf = sp.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Calibri"
        set_ea(run)
    return sp


def add_arrow(slide, x, y, w, h, down=True, color=PRIMARY):
    shape = MSO_SHAPE.DOWN_ARROW if down else MSO_SHAPE.RIGHT_ARROW
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def add_connector(slide, x1, y1, x2, y2, color=PRIMARY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)
    conn.shadow.inherit = False
    return conn


def add_title(slide, text, sub=None):
    add_text(slide, 0.5, 0.24, 12.3, 0.55, [text], size=27, color=PRIMARY, bold=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.55, 0.86, 1.15, 0.045, fill=PRIMARY, line=None)
    if sub:
        add_text(slide, 0.55, 0.94, 11.5, 0.3, [sub], size=10.5, color=MUTED)


def add_footer(slide, page):
    add_text(slide, 0.55, 7.08, 6.0, 0.28, ["AI 能力测评智能体 · A01 数字马力"], size=9, color=MUTED)
    add_text(slide, 12.3, 7.08, 0.85, 0.28, [str(page) + " / 12"], size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, number, title, body, body_size=11.5):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=CARD_BG, line=CARD_LINE)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.12, y + 0.12, 0.32, 0.32, str(number),
              fill=PRIMARY, line=None, size=13, color=WHITE, bold=True)
    add_text(slide, x + 0.12, y + 0.5, w - 0.24, 0.3, [title], size=13, color=DARK, bold=True)
    add_text(slide, x + 0.12, y + 0.82, w - 0.24, h - 0.95, [body], size=body_size, color=MUTED, space_after=2)


def add_mermaid(slide, code, x, y, w, h, label="图：mermaid 代码（复制到 mermaid.live 导出）"):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=CODE_BG, line=CODE_LINE)
    add_text(slide, x + 0.15, y + 0.06, w - 0.3, 0.22, [label], size=9, color=MUTED)
    box = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.3), Inches(w - 0.3), Inches(h - 0.42))
    tf = box.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    first = True
    for line in code.splitlines():
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.0
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9.5)
        run.font.color.rgb = CODE_TEXT
        run.font.name = "Consolas"
        set_ea(run, "等线")
    return box


# ---------------- slide builders ----------------

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.09, fill=PRIMARY, line=None)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0, 7.41, 13.333, 0.09, fill=PRIMARY, line=None)
    add_text(s, 0.8, 2.35, 11.7, 1.1, ["AI 能力测评智能体"], size=46, color=PRIMARY, bold=True,
             align=PP_ALIGN.CENTER)
    add_text(s, 0.8, 3.5, 11.7, 0.6, ["六维能力测评 · 自适应出题 · 大模型自动评分 · 三角色平台"],
             size=18, color=DARK, align=PP_ALIGN.CENTER)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.4, 4.45, 6.53, 1.05, fill=CARD_BG, line=CARD_LINE)
    add_text(s, 3.55, 4.62, 6.23, 0.75,
             ["版本 v1.5.0　题库 301 题　六维模型　单元测试 63 项全绿",
              "A01 数字马力杯 · 技术方案演示"], size=13, color=MUTED, align=PP_ALIGN.CENTER, space_after=4)
    add_footer(s, 1)


def slide_problems(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "要解决的问题", "面向高校与企业 AI 素养测评的现状缺口")
    cards = [
        ("能力标准缺失", "学员不清楚该掌握哪些 AI 能力，缺乏可对标的等级标准"),
        ("自评偏差大", "学员往往高估或低估自身水平，难以获得客观反馈"),
        ("测评手段单一", "传统考试无法考核人机协作、提示词工程、结果甄别等动态技能"),
        ("规模化困难", "人工测评成本高、耗时长，难以在班级或全校范围开展"),
    ]
    for i, (t, b) in enumerate(cards):
        col, row = i % 2, i // 2
        add_card(s, 0.7 + col * 6.15, 1.5 + row * 2.55, 5.9, 2.25, i + 1, t, b)
    add_footer(s, 2)


def slide_dimensions(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "六维能力模型", "每维 50 题，L1~L5 行为锚定分级")
    dims = [
        ("AI 基础认知", "模型原理、能力边界、上下文与幻觉"),
        ("提示词工程", "指令清晰、任务拆解、少样本、格式约束"),
        ("AI 工具使用", "通用大模型、办公插件、数据分析工具选型与组合"),
        ("结果评估与优化", "事实核验、幻觉识别、评分体系与迭代"),
        ("人机协同", "分工、检查点、异常升级、人在回路"),
        ("伦理与合规", "隐私、版权、公平性、可解释性"),
    ]
    for i, (t, b) in enumerate(dims):
        col, row = i % 2, i // 2
        add_card(s, 0.7 + col * 6.15, 1.45 + row * 1.62, 5.9, 1.42, i + 1, t, b, body_size=11)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 6.35, 11.93, 0.5,
              "题型结构：客观 30 + 开放/实操各 4 + 对话/代码/图像/判断各 3 · 答案位置 A/B/C/D 各 45",
              fill=PRIMARY_SOFT, line=None, size=12, color=PRIMARY, bold=True)
    add_footer(s, 3)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "技术架构", "四层分离 + LLM 评分适配层")
    # left: layered diagram with shapes
    lx = 0.7
    layers = [
        ("前端：原生 HTML/CSS/JS 单页", "Canvas 雷达图 · localStorage 状态 · 无构建依赖", PRIMARY_SOFT, PRIMARY),
        ("API 层：FastAPI", "测评/报告 · AI助手/对话 · 题库/复核 · 企业/岗位", CARD_BG, DARK),
        ("领域服务层", "自适应测评引擎 · LLM 评分适配 · 岗位匹配与匿名聚合", CARD_BG, DARK),
        ("数据层：SQLite + JSON 种子", "assessment.db · 题库 301 题 · evidence 证据文件", CARD_BG, DARK),
    ]
    y = 1.35
    for i, (t, b, fill, col) in enumerate(layers):
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, lx, y, 5.7, 1.02, t + "\n" + b,
                  fill=fill, line=CARD_LINE, size=11.5, color=col, bold=(i == 0))
        if i < len(layers) - 1:
            add_arrow(s, lx + 2.7, y + 1.05, 0.3, 0.28)
        y += 1.36
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.75, 2.6, 2.3, 1.5, "LLM 评分/对话\nant-line · 百宝箱\nOpenAI 兼容",
              fill=PRIMARY_SOFT, line=PRIMARY, size=11, color=PRIMARY, bold=True)
    add_connector(s, 6.4, 3.3, 6.75, 3.3)
    # right: notes
    add_text(s, 6.75, 4.35, 5.9, 2.4, [
        "• FastAPI + Pydantic 入参校验，非法输入直接 422",
        "• SQLite 单文件，外键约束，事务自动提交",
        "• 引擎状态快照存 JSON，刷新/重启可续答",
        "• 评分失败显式降级，不冒充模型结果",
    ], size=12, color=DARK, space_after=7)
    add_footer(s, 4)


def slide_engine(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "自适应测评引擎", "分层随机 + 置信度补测 + 末段主观题")
    # flow diagram (shapes)
    y = 1.35
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, y, 2.1, 0.95, "初始 12 题\n六维轮转，难度 1/2",
              fill=PRIMARY_SOFT, line=PRIMARY, size=11, color=PRIMARY, bold=True)
    add_shape(s, MSO_SHAPE.DIAMOND, 3.0, y + 0.12, 1.6, 0.85, "已答 ≥ 12 ?", fill=CARD_BG, line=CARD_LINE, size=10.5, color=DARK)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 4.9, y, 2.1, 0.95, "自适应补测\n置信度/薄弱维度排序", fill=CARD_BG, line=CARD_LINE, size=10.5, color=DARK)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 7.3, y, 2.1, 0.95, "末段主观题\n开放/实操/对话/代码/图像", fill=CARD_BG, line=CARD_LINE, size=10.5, color=DARK)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 9.7, y, 2.1, 0.95, "生成报告\n六维分数 + 雷达图", fill=CARD_BG, line=CARD_LINE, size=10.5, color=DARK)
    add_arrow(s, 2.85, y + 0.36, 0.22, 0.24)
    add_arrow(s, 4.62, y + 0.36, 0.22, 0.24)
    add_arrow(s, 7.02, y + 0.36, 0.22, 0.24)
    add_arrow(s, 9.42, y + 0.36, 0.22, 0.24)
    add_text(s, 3.2, y + 1.1, 4.5, 0.3, ["否 → 继续作答；是 → 进入补测"], size=9.5, color=MUTED)
    # notes
    add_text(s, 0.7, 3.1, 11.9, 1.5, [
        "• 计分：加权正确率 + 证据可靠度，从 50 分先验平滑收敛，单题不跳 0/100",
        "• 同场按题目 ID 去重；候选题过滤由严到宽，题库稀疏时仍可抽到题",
        "• 15 题场末段 3 类主观题，18/25 题场完整纳入 5 类（含代码/图像）",
    ], size=12, space_after=6)
    add_mermaid(s, """flowchart TD
    A[初始 12 题: 六维轮转] --> B{已答 >= 12?}
    B -- 否 --> C[继续作答]
    B -- 是 --> D[按 置信度/分数 排序维度补题]
    D --> E{达到目标题量?}
    E -- 否 --> D
    E -- 是 --> F[末段 5 类主观题]
    F --> G[生成报告]""", 0.7, 4.75, 11.93, 2.05)
    add_footer(s, 5)


def slide_scoring(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "大模型调用与评分", "四通道降级 + 结构化校验 + 人工复核闭环")
    # degradation chain (shapes)
    y = 1.35
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, y, 2.3, 1.05, "主观题作答\n题目 + 答案 + 对话记录", fill=PRIMARY_SOFT, line=PRIMARY, size=11, color=PRIMARY, bold=True)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 3.6, y, 2.9, 1.05, "评分通道\n百宝箱 → ant-line → OpenAI", fill=CARD_BG, line=CARD_LINE, size=11, color=DARK)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 7.1, y, 2.5, 1.05, "失败 → 规则评分\nrubric-fallback", fill=CARD_BG, line=CARD_LINE, size=11, color=DARK)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 10.2, y, 2.4, 1.05, "JSON 校验\n分数钳制 0~满分", fill=CARD_BG, line=CARD_LINE, size=11, color=DARK)
    add_arrow(s, 3.05, y + 0.4, 0.25, 0.24)
    add_arrow(s, 6.55, y + 0.4, 0.25, 0.24)
    add_arrow(s, 9.65, y + 0.4, 0.25, 0.24)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 10.2, 2.85, 2.4, 0.95, "置信度 < 0.75\n强制人工复核", fill=PRIMARY_SOFT, line=PRIMARY, size=11, color=PRIMARY, bold=True)
    add_arrow(s, 11.4, 2.42, 0.24, 0.28)
    add_text(s, 0.7, 2.9, 9.2, 1.5, [
        "• 每次远程调用 60s 超时；失败显式标注 warning 降级，不冒充模型",
        "• AI 助手按轮次回复，历史回放保证上下文连续（多轮剧本已修复并回归）",
        "• 模型返回 response_format=json_object，解析失败同样降级",
    ], size=12, space_after=6)
    add_mermaid(s, """flowchart TD
    A[主观题作答] --> B{通道优先级}
    B -- 百宝箱 --> C[百宝箱 /api/chat]
    B -- ant-line --> D[ant-line /chat/completions]
    B -- OpenAI --> E[Responses API]
    C/D/E -- 失败 --> F[规则评分 rubric-fallback]
    F --> G[JSON校验+分数钳制]
    G --> H{置信度<0.75?}
    H -- 是 --> I[人工复核队列]
    H -- 否 --> J[入库]""", 0.7, 4.75, 11.93, 2.05)
    add_footer(s, 6)


def slide_storage(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "数据存储与安全")
    add_text(s, 0.7, 1.35, 5.6, 3.6, [
        "• SQLite 单文件 data/assessment.db；外键约束、事务自动提交",
        "• A01_DB_PATH 环境变量切换测试库，单元测试不碰真实数据",
        "• 题库 JSON 种子 + 数据库版本留痕（修改/启停记 question_versions）",
        "• 教师/企业口令 PBKDF2 哈希（12 万次迭代），会话令牌 7 天",
        "• 企业端只看匿名画像与匹配分，不接触原始答案",
        "• 学员数据最小化收集，可按测评维度追溯",
    ], size=12.5, space_after=9)
    # storage table
    rows, cols = 5, 3
    table = s.shapes.add_table(rows, cols, Inches(6.7), Inches(1.5), Inches(5.9), Inches(3.6)).table
    data = [
        ("数据", "存储位置", "说明"),
        ("测评记录", "SQLite", "用户/测评/答案/复核/版本"),
        ("题库", "JSON 种子 + 库", "301 题 + 启停版本留痕"),
        ("证据文件", "data/evidence/", "图片/文本，存路径 + SHA256"),
        ("密钥", "环境变量/配置文件", "不写入代码与数据库"),
    ]
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = data[r][c]
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = DARK if r else WHITE
                p.font.name = "Calibri"
                for run in p.runs:
                    set_ea(run)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = PRIMARY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = CARD_BG if r % 2 else WHITE
    add_footer(s, 7)


def slide_demo(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "演示流程", "六步走通完整测评闭环")
    steps = [
        ("选择模式", "首页选择 15/18/25 题，开始测评"),
        ("客观题作答", "选项随机打乱，提交即出判定"),
        ("AI 助手对话", "真实模型，对话计入评分证据"),
        ("生成报告", "雷达图、置信度、错题解析"),
        ("成长档案", "历次测评，点击进入单次报告"),
        ("教师/企业端", "看板画像与岗位投递闭环"),
    ]
    for i, (t, b) in enumerate(steps):
        col, row = i % 3, i // 3
        add_card(s, 0.7 + col * 4.1, 1.5 + row * 2.5, 3.85, 2.2, i + 1, t, b)
    add_footer(s, 8)


def slide_validation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "验证情况", "自动化测试 + 题库结构校验 + 仿真双评")
    metrics = [
        ("63 项", "单元测试全绿", "引擎/评分/鉴权/岗位/选项/多模态回归"),
        ("76 项", "API 冒烟回归通过", "测评全流程 + 管理端 + 企业端"),
        ("301 题", "题库结构校验 valid", "基础认知 51 题，其余五维各 50 题"),
        ("0.932", "仿真模型-人工相关", "100 份样本 / 400 条评分（合成数据）"),
    ]
    for i, (num, t, b) in enumerate(metrics):
        col, row = i % 2, i // 2
        x, y = 0.7 + col * 6.15, 1.5 + row * 2.4
        add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, 5.9, 2.1, fill=CARD_BG, line=CARD_LINE)
        add_text(s, x + 0.2, y + 0.22, 2.2, 0.8, [num], size=30, color=PRIMARY, bold=True)
        add_text(s, x + 2.5, y + 0.32, 3.3, 0.35, [t], size=14, color=DARK, bold=True)
        add_text(s, x + 0.2, y + 1.25, 5.5, 0.7, [b], size=11, color=MUTED)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 6.4, 11.93, 0.5,
              "说明：等级阈值由仿真样本校准；真实学员试测与教师双评待补", fill=PRIMARY_SOFT, line=None, size=12, color=PRIMARY, bold=True)
    add_footer(s, 9)


def slide_future(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "后续工作")
    cards = [
        ("真实验证", "30 人真实试测 + 教师双评，校准等级阈值与题目难度"),
        ("代码沙箱", "代码任务接入真实执行环境；图像题库替换真实对比素材"),
        ("数据层演进", "迁移 PostgreSQL，支持多租户与更大并发"),
        ("生态对接", "LMS/证书体系：能力画像标准 JSON、授权投递与申诉"),
    ]
    for i, (t, b) in enumerate(cards):
        col, row = i % 2, i // 2
        add_card(s, 0.7 + col * 6.15, 1.7 + row * 2.4, 5.9, 2.1, i + 1, t, b)
    add_footer(s, 10)


def slide_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "总结")
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 0.7, 1.5, 5.85, 4.4, fill=CARD_BG, line=CARD_LINE)
    add_text(s, 0.95, 1.72, 5.35, 0.4, ["已完成"], size=16, color=PRIMARY, bold=True)
    add_text(s, 0.95, 2.2, 5.35, 3.4, [
        "• 六维模型 + 301 题库 + 自适应引擎",
        "• 真实大模型评分（ant-line/百宝箱）",
        "• 报告、成长档案、教师看板",
        "• 企业岗位匹配与匿名画像闭环",
        "• 人工双评、争议裁决、版本留痕",
    ], size=13, space_after=8)
    add_shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 6.8, 1.5, 5.85, 4.4, fill=PRIMARY_SOFT, line=PRIMARY)
    add_text(s, 7.05, 1.72, 5.35, 0.4, ["下一步"], size=16, color=PRIMARY, bold=True)
    add_text(s, 7.05, 2.2, 5.35, 3.4, [
        "• 从功能开发转向真实验证",
        "• 真实试测 + 教师双评数据",
        "• 演示视频与稳定在线环境",
        "• 生态对接（LMS/证书/招聘）",
    ], size=13, space_after=8)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0.55, 6.35, 12.23, 0.5, "评分三层降级、人工复核闭环、数据脱敏与版本留痕 —— 支撑真实效度验证",
              fill=PRIMARY, line=None, size=13, color=WHITE, bold=True)
    add_footer(s, 11)


def slide_end(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_shape(s, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.09, fill=PRIMARY, line=None)
    add_shape(s, MSO_SHAPE.RECTANGLE, 0, 7.41, 13.333, 0.09, fill=PRIMARY, line=None)
    add_text(s, 0.8, 3.0, 11.7, 0.9, ["谢谢"], size=40, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, 0.8, 4.1, 11.7, 1.2, [
        "演示环境：本地一键启动（python run.py），局域网 / 临时隧道可访问",
        "评分通道：ant-line · 百宝箱 · OpenAI 兼容 · 本地规则四层降级",
    ], size=14, color=MUTED, align=PP_ALIGN.CENTER, space_after=8)
    add_footer(s, 12)


def build():
    prs = Presentation()
    prs.slide_width = Emu(EMU_W)
    prs.slide_height = Emu(EMU_H)
    for fn in (slide_cover, slide_problems, slide_dimensions, slide_architecture,
               slide_engine, slide_scoring, slide_storage, slide_demo,
               slide_validation, slide_future, slide_summary, slide_end):
        fn(prs)
    prs.save(OUT)
    print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    build()
